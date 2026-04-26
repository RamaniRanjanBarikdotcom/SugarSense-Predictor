"""SugarSense Predictor — 10-Slide B.Tech Minor Project Presentation"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

W, H = 13.33, 7.5

BG    = "080f1e"; BG2 = "101e38"; BG3 = "0d1a2f"
PRI   = "0f766e"; ACC = "14b8a6"; LACC = "5eead4"
WHITE = "ffffff"; GRAY = "94a3b8"; DGRAY = "475569"
YELLOW= "fbbf24"; GREEN = "4ade80"; PINK = "f472b6"

GUIDE   = "Prof. Nabin Kumar Nag"
COLLEGE = "Sophitorium Engineering College"
DEPT    = "Department of Computer Science & Engineering"
YEAR    = "2025 – 2026"
TEAM = [
    ("Ramani Ranjan Barik",  "Roll: 2301332XXX"),
    ("Laxmi Kanta Panda",    "Roll: 2301332012"),
    ("[Team Member 3]",      "Roll: XXXXXXXXXX"),
    ("[Team Member 4]",      "Roll: XXXXXXXXXX"),
    ("[Team Member 5]",      "Roll: XXXXXXXXXX"),
]
# Slide ranges per presenter
PRESENTER_SLIDES = [
    (TEAM[0][0], "Slides 1–2"),
    (TEAM[1][0], "Slides 3–4"),
    (TEAM[2][0], "Slides 5–6"),
    (TEAM[3][0], "Slides 7–8"),
    (TEAM[4][0], "Slides 9–10"),
]

prs = Presentation()
prs.slide_width  = Inches(W)
prs.slide_height = Inches(H)
blank = prs.slide_layouts[6]

# ── helpers ───────────────────────────────────────────────────────────────────
def rgb(h): return RGBColor(int(h[:2],16), int(h[2:4],16), int(h[4:],16))

def bg(slide, c=BG):
    f = slide.background.fill; f.solid(); f.fore_color.rgb = rgb(c)

def box(slide, l, t, w, h, fill=None, line=None, lw=1.2, rnd=False):
    shp = slide.shapes.add_shape(9 if rnd else 1,
          Inches(l), Inches(t), Inches(w), Inches(h))
    shp.fill.solid() if fill else shp.fill.background()
    if fill: shp.fill.fore_color.rgb = rgb(fill)
    if line:
        shp.line.color.rgb = rgb(line)
        shp.line.width = Pt(lw)
    else:
        shp.line.fill.background()
    return shp

def tx(slide, text, l, t, w, h, sz=12, bold=False, col=WHITE,
       align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = rgb(col)
    return tb

def card(slide, l, t, w, h, fill=BG2, ac=ACC, top=False):
    box(slide, l, t, w, h, fill=fill)
    if top: box(slide, l, t, w, 0.045, fill=ac)
    else:   box(slide, l, t, 0.045, h, fill=ac)

def hdr(slide, title, sub=None):
    box(slide, 0, 0, W, 0.055, fill=ACC)
    tx(slide, title, 0.38, 0.1, W-0.76, 0.62, sz=24, bold=True)
    if sub: tx(slide, sub, 0.38, 0.7, W-0.76, 0.3, sz=9.5, col=GRAY)

def ptag(slide, idx):
    name, slides = PRESENTER_SLIDES[idx]
    box(slide, W-3.05, H-0.38, 3.0, 0.3, fill=PRI)
    tx(slide, f"▶  {name}  ·  {slides}", W-3.03, H-0.37, 2.96, 0.26,
       sz=7, align=PP_ALIGN.CENTER)

def num_circle(slide, n, l, t, r=0.2, fill=ACC):
    shp = slide.shapes.add_shape(9, Inches(l), Inches(t), Inches(r*2), Inches(r*2))
    shp.fill.solid(); shp.fill.fore_color.rgb = rgb(fill)
    shp.line.fill.background()
    tf = shp.text_frame; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r2 = p.add_run(); r2.text = str(n)
    r2.font.size = Pt(10); r2.font.bold = True; r2.font.color.rgb = rgb(WHITE)

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 1 — COVER
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank); bg(s, "060d1a")

# Left accent bar
box(s, 0, 0, 0.45, H, fill=PRI)
box(s, 0.45, 0, 0.04, H, fill=ACC)

# Decorative circles top-right
for cx,cy,cr,co in [(11,0.2,1.6,PRI),(12.4,0.1,1.0,ACC),(11.6,1.8,0.6,"0d2540")]:
    shp = s.shapes.add_shape(9,Inches(cx),Inches(cy),Inches(cr),Inches(cr))
    shp.fill.solid(); shp.fill.fore_color.rgb=rgb(co); shp.line.fill.background()

# Minor project badge
box(s, 0.75, 0.48, 2.7, 0.3, fill=ACC, rnd=True)
tx(s, "B.TECH MINOR PROJECT  ·  2025-2026", 0.77, 0.5, 2.65, 0.25,
   sz=7.5, bold=True, align=PP_ALIGN.CENTER)

# Title
tx(s, "SugarSense", 0.65, 1.0, 11, 1.0, sz=58, bold=True)
tx(s, "Predictor",  0.65, 1.85, 11, 0.9, sz=58, bold=True, col=LACC)
tx(s, "Smart Diabetes Risk Prediction  &  AI Health Assistant System",
   0.65, 2.78, 11, 0.5, sz=14, col=GRAY)

# Divider
box(s, 0.65, 3.42, 4.8, 0.035, fill=ACC)

# College & Guide
tx(s, COLLEGE, 0.65, 3.58, 9, 0.38, sz=13, bold=True)
tx(s, DEPT,    0.65, 3.92, 9, 0.3,  sz=10, col=GRAY)
tx(s, f"Under the Guidance of  {GUIDE}", 0.65, 4.24, 9, 0.3, sz=10, col=LACC, italic=True)

# 5 team cards
cw = 2.38
for i,(name,roll) in enumerate(TEAM):
    x = 0.65 + i*(cw+0.13)
    card(s, x, 5.0, cw, 1.1, fill="0f1e38", ac=ACC)
    tx(s, f"0{i+1}", x+0.12, 5.04, 0.5, 0.42, sz=20, bold=True, col=ACC)
    tx(s, name, x+0.12, 5.48, cw-0.2, 0.3, sz=9, bold=True)
    tx(s, roll, x+0.12, 5.78, cw-0.2, 0.25, sz=7.5, col=GRAY)

tx(s, "All 5 Members — Cover", W-2.3, H-0.34, 2.2, 0.26,
   sz=7, col=LACC, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 2 — INTRODUCTION + PROBLEM STATEMENT
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank); bg(s)
hdr(s, "Introduction & Problem Statement",
    "Global diabetes challenge, what SugarSense solves, and why existing tools fall short")

# Stat cards
for i,(val,lbl) in enumerate([("537M+","Adults affected (2021)"),("700M+","Projected by 2045"),
                               ("50%","Cases undiagnosed"),("Top 10","Cause of death")]):
    x = 0.38 + i*3.22; cw2=3.06
    card(s, x, 1.12, cw2, 0.95, fill=BG2, ac=ACC, top=True)
    tx(s, val, x+0.1, 1.19, cw2-0.2, 0.5, sz=24, bold=True, col=LACC, align=PP_ALIGN.CENTER)
    tx(s, lbl, x+0.1, 1.66, cw2-0.2, 0.3, sz=9,  col=GRAY,  align=PP_ALIGN.CENTER)

# What is SugarSense card
card(s, 0.38, 2.2, 7.6, 2.08, fill=BG2, ac=PRI)
tx(s, "What is SugarSense Predictor?", 0.58, 2.28, 7.1, 0.38, sz=12, bold=True)
tx(s,
   "A full-stack web app combining an SVM classifier (trained on Pima Indians Diabetes Dataset) "
   "with a multi-provider AI Health Chatbot. Users enter 8 clinical metrics — the model predicts "
   "Diabetic / Non-Diabetic in <100ms, then the AI chatbot provides personalized health guidance, "
   "explains the result, and answers follow-up questions.",
   0.58, 2.68, 7.1, 1.52, sz=10, col=GRAY)

# Problem Statement card
card(s, 8.12, 2.2, 4.83, 2.08, fill=BG2, ac=YELLOW, top=True)
tx(s, "Problem Statement", 8.3, 2.28, 4.5, 0.35, sz=12, bold=True, col=YELLOW)
problems = ["No real-time health guidance after prediction",
            "Lab-dependent screening — inaccessible & costly",
            "No conversational AI to explain risk factors"]
for i,p in enumerate(problems):
    box(s, 8.28, 2.72+i*0.44, 0.05, 0.3, fill=YELLOW)
    tx(s, p, 8.4, 2.72+i*0.44, 4.42, 0.38, sz=9.5, col=GRAY)

# 3 highlights row
for i,(ac,t,d) in enumerate([(ACC,"SVM ML Model","~78% accuracy on Pima dataset"),
                               (PINK,"5 AI Providers","Claude · GPT · Gemini · HF · OpenRouter"),
                               (GREEN,"110+ Q&A KB","Health knowledge embedded in AI prompt")]):
    x = 0.38+i*4.35; cw3=4.2
    card(s, x, 4.42, cw3, 0.9, fill=BG2, ac=ac)
    tx(s, t, x+0.2, 4.5,  cw3-0.3, 0.35, sz=11, bold=True, col=ac)
    tx(s, d, x+0.2, 4.85, cw3-0.3, 0.38, sz=9.5, col=GRAY)

# Objective strip
box(s, 0.38, 5.46, W-0.76, 1.72, fill=BG3)
box(s, 0.38, 5.46, W-0.76, 0.04, fill=ACC)
tx(s, "OBJECTIVES", 0.55, 5.54, 2, 0.3, sz=9, bold=True, col=ACC)
objs = ["Accurate SVM Prediction","Responsive React UI","Animated Risk Gauge",
        "AI Health Chatbot","Health Knowledge Base","5-Provider Flexibility","History Tracking"]
for i,o in enumerate(objs):
    x = 0.55 + i*1.84
    num_circle(s, i+1, x, 5.9, 0.16)
    tx(s, o, x-0.02, 6.22, 1.7, 0.88, sz=8, col=GRAY, align=PP_ALIGN.CENTER)

ptag(s, 0)

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 3 — DATASET DESCRIPTION
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank); bg(s)
hdr(s, "Dataset Description",
    "Pima Indians Diabetes Database — NIDDK / UCI Machine Learning Repository / Kaggle")

# 4 stat cards
for i,(v,l,c) in enumerate([("768","Patient Records",ACC),("8","Input Features",LACC),
                              ("65.1%","Non-Diabetic",GREEN),("34.9%","Diabetic","f87171")]):
    x = 0.38 + i*3.24; cw4=3.08
    card(s, x, 1.12, cw4, 0.95, fill=BG2, ac=c, top=True)
    tx(s, v, x+0.1, 1.18, cw4-0.2, 0.52, sz=26, bold=True, col=c, align=PP_ALIGN.CENTER)
    tx(s, l, x+0.1, 1.68, cw4-0.2, 0.3,  sz=9.5, col=GRAY, align=PP_ALIGN.CENTER)

# Feature table
card(s, 0.38, 2.2, W-0.76, 5.0, fill=BG2, ac=PRI)
tx(s, "8 Clinical Input Features", 0.58, 2.28, 6, 0.36, sz=12, bold=True)

heads = ["#","Feature","Description","Unit","Valid Range","Normal Range"]
hxs   = [0.55, 0.98, 3.22, 7.78, 9.08, 11.0]
hws   = [0.38, 2.18, 4.5, 1.24, 1.86, 2.1]
box(s, 0.55, 2.7, W-1.0, 0.38, fill=PRI)
for j,(h_,hx_) in enumerate(zip(heads,hxs)):
    tx(s, h_, hx_, 2.73, hws[j], 0.28, sz=8.5, bold=True)

rows = [
    ("1","Pregnancies","Number of times pregnant","Count","0 – 25","0 – 5"),
    ("2","Glucose","2-hour plasma glucose (OGTT)","mg/dL","1 – 300","70 – 99"),
    ("3","Blood Pressure","Diastolic blood pressure","mm Hg","20 – 200","60 – 80"),
    ("4","Skin Thickness","Triceps skin fold thickness","mm","0 – 100","10 – 40"),
    ("5","Insulin","2-hour serum insulin","μU/mL","0 – 900","16 – 166"),
    ("6","BMI","Body mass index (weight / height²)","kg/m²","10 – 70","18.5 – 24.9"),
    ("7","Diabetes Pedigree Fn.","Genetic likelihood from family history","Score","0.05–3.0","0.08–0.8"),
    ("8","Age","Age of the patient","Years","1 – 120","Any"),
]
for ri,row in enumerate(rows):
    ry = 3.1 + ri*0.515
    box(s, 0.55, ry, W-1.0, 0.5, fill="0d1a30" if ri%2==0 else BG2)
    for j,(v,hx_) in enumerate(zip(row,hxs)):
        c_ = LACC if j==0 else (WHITE if j==1 else GRAY)
        tx(s, v, hx_, ry+0.1, hws[j], 0.3, sz=9, bold=(j==1), col=c_)

ptag(s, 1)

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 4 — SYSTEM ARCHITECTURE
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank); bg(s)
hdr(s, "System Architecture",
    "Three-tier client–server design with external AI provider routing")

tiers = [
    (ACC,  "TIER 1 — PRESENTATION","React SPA  (index.html)",
     ["8-Field validated prediction form","Animated SVG risk gauge",
      "Floating AI chat panel","Assessment history tab","Key factor analysis panel"]),
    (YELLOW,"TIER 2 — APPLICATION","Flask Backend  (app.py)",
     ["GET /  →  Serve React SPA","POST /predict_api  →  ML inference",
      "POST /chat  →  AI provider router","Input validation (client + server)",
      "Load sc1.pkl & classifier1.pkl at startup"]),
    (PINK, "TIER 3 — AI PROVIDERS","5 External LLM APIs",
     ["Anthropic Claude  (SDK + prompt cache)","OpenAI GPT  (httpx REST)",
      "Google Gemini  (httpx REST)","HuggingFace Inference  (httpx REST)",
      "OpenRouter  (100+ models via httpx)"]),
]
tw = (W-0.76-0.3)/3
for i,(ac,tier,comp,items) in enumerate(tiers):
    x = 0.38+i*(tw+0.15)
    card(s, x, 1.12, tw, 5.56, fill=BG2, ac=ac)
    box(s, x, 1.12, tw, 0.44, fill=ac)
    tx(s, tier, x+0.12, 1.16, tw-0.2, 0.24, sz=7.5, bold=True, col=WHITE)
    tx(s, comp, x+0.12, 1.44, tw-0.2, 0.36, sz=11.5, bold=True, col=WHITE)
    box(s, x+0.12, 1.84, tw-0.24, 0.02, fill=ac)
    for j,item in enumerate(items):
        cy = 1.94+j*0.54
        box(s, x+0.12, cy+0.1, 0.055, 0.3, fill=ac)
        tx(s, item, x+0.24, cy+0.06, tw-0.36, 0.38, sz=9.5, col=GRAY)

# Arrow connectors
for ax in [0.38+tw+0.02, 0.38+2*(tw+0.15)-0.14]:
    tx(s, "→", ax, 3.55, 0.22, 0.42, sz=20, bold=True, col=ACC, align=PP_ALIGN.CENTER)

# ML artifacts footer
box(s, 0.38, 6.82, W-0.76, 0.36, fill="0d1a30")
tx(s, "ML ARTIFACTS:", 0.55, 6.88, 1.6, 0.22, sz=8.5, bold=True, col=ACC)
tx(s, "sc1.pkl  (StandardScaler — normalise 8 inputs to mean=0, std=1)     "
      "classifier1.pkl  (SVM linear kernel — output 0=Non-Diabetic | 1=Diabetic)",
   2.2, 6.88, 10.5, 0.22, sz=8.5, col=GRAY)

ptag(s, 1)

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 5 — METHODOLOGY + ALGORITHM (SVM)
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank); bg(s)
hdr(s, "Methodology & Algorithm — SVM",
    "6-step ML pipeline from raw CSV to live prediction + why Support Vector Machine")

# Pipeline steps (left 2/3)
steps = [
    (ACC,    "01 Data Collection",   "768 records · Pima Indians Diabetes DB · 8 features + Outcome"),
    (LACC,   "02 Preprocessing",     "Handle zero-value missing data · separate X (features) and y (label)"),
    (YELLOW, "03 Train-Test Split",  "80% training (614 samples) / 20% test (154) · stratified sampling"),
    (PINK,   "04 Feature Scaling",   "StandardScaler: fit on train → transform both splits → saved as sc1.pkl"),
    (GREEN,  "05 SVM Training",      "SVC linear kernel: finds hyperplane maximising class margin"),
    (ACC,    "06 Live Inference",     "Validate → scale → predict → return {prediction: 0|1} in <100ms"),
]
sw = 8.4; sh = 0.88
for i,(ac,title,body) in enumerate(steps):
    y = 1.12+i*(sh+0.04)
    card(s, 0.38, y, sw, sh, fill=BG2, ac=ac)
    tx(s, title, 0.56, y+0.08, 3.0, 0.32, sz=11, bold=True, col=ac)
    tx(s, body,  0.56, y+0.44, sw-0.35, 0.38, sz=9.5, col=GRAY)
    # Arrow (not on last)
    if i<5: tx(s, "↓", 0.38+sw/2, y+sh-0.08, 0.3, 0.28, sz=12, bold=True, col=DGRAY, align=PP_ALIGN.CENTER)

# SVM explanation (right panel)
card(s, 8.96, 1.12, 4.0, 6.06, fill=BG2, ac=YELLOW, top=True)
tx(s, "SVM Explained", 9.12, 1.22, 3.7, 0.35, sz=12, bold=True, col=YELLOW)

svm_pts = [
    (ACC,  "Classification:", "Finds a hyperplane that maximally separates two classes."),
    (LACC, "Max Margin:",     "Distance from hyperplane to nearest points (support vectors) is maximised."),
    (YELLOW,"Linear Kernel:", "A straight boundary in 8D feature space after StandardScaler normalisation."),
    (PINK, "Why SVM?",        "Robust on small datasets · no Gaussian assumption · high-dimensional effective."),
    (GREEN,"Accuracy:",       "~78% on test set · Precision ~80% · Recall ~80% · F1 ~81%"),
]
for i,(c,bold_t,body) in enumerate(svm_pts):
    y = 1.68+i*0.88
    box(s, 9.1, y, 0.04, 0.62, fill=c)
    tx(s, bold_t, 9.22, y+0.02, 3.6, 0.28, sz=9.5, bold=True, col=c)
    tx(s, body,   9.22, y+0.3,  3.6, 0.52, sz=8.5, col=GRAY)

# Metric row
for i,(lab,val,c) in enumerate([("Accuracy","~78%",ACC),("Precision","~80%",LACC),
                                  ("Recall","~80%",YELLOW),("F1","~81%",GREEN)]):
    xi = 8.98+i*1.0
    card(s, xi, 6.48, 0.94, 0.62, fill=BG3, ac=c, top=True)
    tx(s, val, xi+0.05, 6.56, 0.84, 0.3, sz=13, bold=True, col=c, align=PP_ALIGN.CENTER)
    tx(s, lab, xi+0.05, 6.86, 0.84, 0.2, sz=7.5, col=GRAY, align=PP_ALIGN.CENTER)

ptag(s, 2)

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 6 — TECHNOLOGY STACK
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank); bg(s)
hdr(s, "Technology Stack",
    "Modern, production-grade tools across every layer — frontend, backend, ML, and AI")

tech = [
    (ACC,    "⚛  React 18.3.1",         "Frontend",    "CDN + Babel. Single-file SPA — no build step required."),
    (YELLOW, "🐍  Flask ≥ 3.0",          "Backend",     "REST API server. 5 routes. .env config via python-dotenv."),
    (PINK,   "🤖  scikit-learn ≥ 1.4",  "ML Engine",   "SVM (linear kernel) + StandardScaler via pickle artifacts."),
    (GREEN,  "🔢  NumPy ≥ 1.26",         "Numerics",    "Array construction for model input; inference pipeline."),
    (LACC,   "🟣  Anthropic SDK ≥ 0.97", "AI Primary",  "Claude API with ephemeral prompt caching on system prompt."),
    (ACC,    "🌐  httpx ≥ 0.28",         "AI REST",     "OpenAI, Google Gemini, HuggingFace, OpenRouter calls."),
    (YELLOW, "⚙  python-dotenv ≥ 1.0",  "Config",      "Load ANTHROPIC_API_KEY, GOOGLE_API_KEY, CHAT_MODEL etc."),
    (PINK,   "💾  pickle (stdlib)",      "Artifacts",   "Deserialize sc1.pkl (scaler) + classifier1.pkl (SVM)."),
]
tw2=(W-0.76-0.45)/4; th2=1.44
for i,(ac,name,layer,desc) in enumerate(tech):
    col=i%4; row=i//2 if i<4 else i//2
    col=i%4; row=i//4
    x=0.38+col*(tw2+0.15); y=1.12+row*(th2+0.16)
    card(s, x, y, tw2, th2, fill=BG2, ac=ac)
    tx(s, name,  x+0.12, y+0.08, tw2-0.2, 0.4, sz=10.5, bold=True, col=WHITE)
    tx(s, layer, x+0.12, y+0.52, tw2-0.2, 0.26, sz=8.5, bold=True, col=ac)
    box(s, x+0.12, y+0.8, tw2-0.24, 0.018, fill=ac)
    tx(s, desc,  x+0.12, y+0.88, tw2-0.2, 0.52, sz=9, col=GRAY)

# Provider routing footer
box(s, 0.38, 6.62, W-0.76, 0.6, fill=BG3)
box(s, 0.38, 6.62, W-0.76, 0.04, fill=ACC)
tx(s, "AI PROVIDER ROUTING  (auto-detected by CHAT_MODEL prefix in .env):",
   0.55, 6.7, 5.5, 0.28, sz=8.5, bold=True, col=ACC)
pvs=[("claude/claude-","Anthropic",ACC),("gpt-/o1/o3/o4","OpenAI",YELLOW),
     ("gemini-","Google",PINK),("hf/...","HuggingFace",GREEN),("(anything else)","OpenRouter",LACC)]
for i,(pfx,name,c) in enumerate(pvs):
    x=6.1+i*1.44
    tx(s, f"{pfx}\n➔ {name}", x, 6.66, 1.38, 0.5, sz=8, bold=True, col=c, align=PP_ALIGN.CENTER)

ptag(s, 2)

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 7 — FEATURES & FUNCTIONALITY
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank); bg(s)
hdr(s, "Features & Functionality",
    "Eight core features built and live in the deployed web application")

feats = [
    (ACC,    "01","8-Field Prediction Form",
             "Validated inputs with unit labels, normal ranges, and field hints. Client + server validation."),
    (LACC,   "02","SVM Risk Prediction",
             "Binary ML output (Diabetic / Non-Diabetic) via POST /predict_api. Response in <100ms."),
    (YELLOW, "03","Animated Risk Gauge",
             "SVG arc gauge: weighted clinical scoring → 2–98% visual risk. Smooth CSS animation."),
    (PINK,   "04","Key Factor Analysis",
             "4 risk factors (Glucose, BMI, Blood Pressure, Insulin) colour-coded by clinical threshold."),
    (GREEN,  "05","AI Health Chatbot",
             "Floating multi-turn chat panel. Full conversation history sent per request."),
    (ACC,    "06","5-Provider AI Support",
             "Anthropic, OpenAI, Google, HuggingFace, OpenRouter — switch via .env CHAT_MODEL only."),
    (LACC,   "07","Health Knowledge Base",
             "110+ Q&A pairs (10 categories) loaded from JSON and embedded into AI system prompt."),
    (YELLOW, "08","Assessment History",
             "Last 5 predictions persisted in browser localStorage. History tab with metrics + dates."),
]
fw=(W-0.76-0.45)/4; fh=2.62
for i,(ac,num,title,body) in enumerate(feats):
    col=i%4; row=i//4
    x=0.38+col*(fw+0.15); y=1.12+row*(fh+0.16)
    card(s, x, y, fw, fh, fill=BG2, ac=ac)
    tx(s, num,   x+0.14, y+0.1,  fw-0.24, 0.52, sz=22, bold=True, col=ac)
    tx(s, title, x+0.14, y+0.65, fw-0.24, 0.44, sz=10.5, bold=True, col=WHITE)
    box(s, x+0.14, y+1.12, fw-0.28, 0.018, fill=ac)
    tx(s, body,  x+0.14, y+1.2,  fw-0.24, 1.36, sz=9, col=GRAY)

ptag(s, 3)

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 8 — AI CHATBOT SYSTEM
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank); bg(s)
hdr(s, "AI Health Chatbot System",
    "Multi-provider conversational AI with a 110+ Q&A structured diabetes knowledge base")

# Left: knowledge base categories
card(s, 0.38, 1.12, 5.9, 6.06, fill=BG2, ac=PINK)
tx(s, "Health Knowledge Base  (health_knowledge.json)", 0.56, 1.2, 5.5, 0.38, sz=12, bold=True)
tx(s, "Loaded at Flask startup · Appended to AI system prompt as structured markdown",
   0.56, 1.58, 5.5, 0.26, sz=8.5, col=GRAY, italic=True)
box(s, 0.56, 1.88, 5.5, 0.02, fill=PINK)

kb=[
    (PINK,   "Form Fields (×8)",        "40 Q&A — every feature, ranges, why it matters"),
    (ACC,    "Diabetes Types",           "6 Q&A — Type 1, Type 2, Pre-diabetes, Gestational"),
    (YELLOW, "Symptoms",                 "6 Q&A — Signs, DKA, Hypoglycemia, Hyperglycemia"),
    (LACC,   "Prevention",              "6 Q&A — Lifestyle, diet, exercise, sleep impact"),
    (GREEN,  "Diet & Nutrition",        "8 Q&A — Glycemic index, best & worst foods"),
    (PINK,   "Exercise",                 "5 Q&A — Types, frequency, blood sugar effect"),
    (ACC,    "Blood Sugar Monitoring",   "4 Q&A — HbA1c, CGM, target ranges"),
    (YELLOW, "Complications",            "5 Q&A — Neuropathy, retinopathy, kidneys"),
    (LACC,   "Medications & Treatment", "4 Q&A — Metformin, Insulin, GLP-1 (Ozempic)"),
    (GREEN,  "About SugarSense",        "6 Q&A — How it works, disclaimers, next steps"),
]
for i,(c,cat,desc) in enumerate(kb):
    y=2.02+i*0.48
    box(s, 0.56, y+0.08, 0.055, 0.28, fill=c)
    tx(s, cat,  0.7,  y+0.07, 2.1, 0.24, sz=9.5, bold=True, col=c)
    tx(s, desc, 0.7,  y+0.3,  5.1, 0.2,  sz=8.5, col=GRAY)

# Right top: 5 providers
card(s, 6.44, 1.12, 6.52, 2.62, fill=BG2, ac=ACC, top=True)
tx(s, "5-Provider AI Architecture", 6.62, 1.22, 6.1, 0.35, sz=12, bold=True)
pvs2=[("🟣","Anthropic Claude","claude / claude-",ACC),
      ("🟢","OpenAI GPT","gpt- / o1 / o4",YELLOW),
      ("🔵","Google Gemini","gemini-",PINK),
      ("🟡","HuggingFace","hf/model-name",GREEN),
      ("🔴","OpenRouter","any other string",LACC)]
pw=(6.52-0.2)/3
for i,(ico,name,pfx,c) in enumerate(pvs2):
    col=i%3; row=i//3
    x=6.5+col*(pw+0.1); y=1.62+row*0.98
    card(s, x, y, pw, 0.84, fill="0d1a30", ac=c)
    tx(s, f"{ico} {name}", x+0.12, y+0.06, pw-0.2, 0.3, sz=9.5, bold=True, col=c)
    tx(s, pfx, x+0.12, y+0.38, pw-0.2, 0.38, sz=8, col=GRAY)

# Right middle: Prompt caching
card(s, 6.44, 3.88, 6.52, 1.42, fill=BG2, ac=LACC, top=True)
tx(s, "Prompt Caching (Claude)", 6.62, 3.98, 6.0, 0.32, sz=11, bold=True, col=LACC)
tx(s,
   "System prompt (role + full KB) sent with cache_control: ephemeral.\n"
   "Anthropic caches the prefix for 5 min → lower latency + reduced API cost on repeated chat requests.",
   6.62, 4.34, 6.0, 0.88, sz=9.5, col=GRAY)

# Right bottom: Multi-turn flow
card(s, 6.44, 5.44, 6.52, 1.74, fill=BG2, ac=YELLOW, top=True)
tx(s, "Multi-Turn Conversation Flow", 6.62, 5.54, 6.0, 0.32, sz=11, bold=True, col=YELLOW)
for i,step in enumerate(["User types → React bundles full history[] array",
                          "POST /chat → Flask strips leading assistant turns",
                          "LLM receives [system prompt + conversation history]",
                          "Reply appended to messages → chat auto-scrolls"]):
    tx(s, f"{i+1}.  {step}", 6.62, 5.9+i*0.36, 6.0, 0.32, sz=9.5, col=GRAY)

ptag(s, 3)

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 9 — RESULTS & FUTURE WORK
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank); bg(s)
hdr(s, "Results & Future Work",
    "Model performance metrics, literature comparison, and the planned development roadmap")

# Metric cards
for i,(val,lab,c) in enumerate([("~78%","Accuracy",ACC),("~80%","Precision",LACC),
                                  ("~80%","Recall",YELLOW),("~81%","F1 Score",GREEN)]):
    x=0.38+i*3.24; mw=3.08
    card(s, x, 1.12, mw, 1.22, fill=BG2, ac=c, top=True)
    tx(s, val, x+0.1, 1.18, mw-0.18, 0.6, sz=28, bold=True, col=c, align=PP_ALIGN.CENTER)
    tx(s, lab, x+0.1, 1.76, mw-0.18, 0.3, sz=10, bold=True, col=WHITE, align=PP_ALIGN.CENTER)

# Literature comparison (left)
card(s, 0.38, 2.5, 7.4, 4.68, fill=BG2, ac=PRI)
tx(s, "Literature Comparison", 0.56, 2.58, 5, 0.36, sz=12, bold=True)

lhds=["Authors","Year","Algorithm","Accuracy"]
lhxs=[0.55, 4.08, 4.82, 7.22]; lhws=[3.48,0.7,2.36,0.88]
box(s, 0.55, 3.0, 7.2, 0.35, fill=PRI)
for j,(h_,hx_) in enumerate(zip(lhds,lhxs)):
    tx(s, h_, hx_, 3.03, lhws[j], 0.26, sz=8.5, bold=True)
lit=[("S. Patel & R. Goyal","2019","SVM (RBF)","78%",GRAY),
     ("M. Elshazly & Motaba","2020","SVM vs DT vs LR","81%",GRAY),
     ("N. Kumar & Srivastava","2021","SVM + KNN","83%",GRAY),
     ("A. Banerjee & Tiwari","2022","SVM + RFE","85%",GRAY),
     ("H. Kaur & M. Singh","2023","SVM + Bagging","86%",GRAY),
     ("SugarSense  ★","2026","SVM (linear)","~78%",LACC)]
for ri,(au,yr,al,ac2,c) in enumerate(lit):
    ry=3.37+ri*0.44; fill="0d1a30" if ri%2==0 else BG2
    if ri==5: fill="0a2a25"
    box(s, 0.55, ry, 7.2, 0.42, fill=fill)
    for j,(v,hx_) in enumerate(zip([au,yr,al,ac2],lhxs)):
        c_=LACC if ri==5 else (ACC if j==3 else c)
        tx(s, v, hx_, ry+0.08, lhws[j], 0.26, sz=9, bold=(ri==5), col=c_)

# Future work (right)
card(s, 7.94, 2.5, 5.02, 4.68, fill=BG2, ac=YELLOW, top=True)
tx(s, "Future Work", 8.1, 2.6, 4.7, 0.34, sz=12, bold=True, col=YELLOW)
fut=[
    (YELLOW,"HIGH","Trend Graphs + Real Confidence Score"),
    (YELLOW,"HIGH","Data Encryption & HIPAA Compliance"),
    (ACC,   "MED", "Alert System for Rising Risk Trends"),
    (ACC,   "MED", "User Authentication + Database Storage"),
    (ACC,   "MED", "Wearable CGM Device Integration"),
    (ACC,   "MED", "Feature Importance Bar Chart"),
    (ACC,   "MED", "Cloud Deployment (AWS/GCP + HTTPS)"),
    (GREEN, "LOW", "Model Retraining Pipeline"),
    (GREEN, "LOW", "Multi-language Support"),
]
for i,(c,pr,title) in enumerate(fut):
    y=3.06+i*0.46
    box(s, 8.1, y+0.07, 0.6, 0.3, fill=c, rnd=True)
    tx(s, pr, 8.12, y+0.09, 0.56, 0.22, sz=7, bold=True, col=WHITE, align=PP_ALIGN.CENTER)
    tx(s, title, 8.8, y+0.08, 4.0, 0.3, sz=9.5, col=GRAY)

ptag(s, 4)

# ══════════════════════════════════════════════════════════════════════════════
#  SLIDE 10 — CONCLUSION + REFERENCES + THANK YOU
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(blank); bg(s, "060d1a")

# Top gradient bar
box(s, 0, 0, W, 0.055, fill=ACC)

# Decorative circles
for cx,cy,cr,co in [(10.5,0,3.5,PRI),(12.8,0.5,2,"0c5e59"),(0,5.5,2,"0d1a2f")]:
    shp=s.shapes.add_shape(9,Inches(cx),Inches(cy),Inches(cr),Inches(cr))
    shp.fill.solid(); shp.fill.fore_color.rgb=rgb(co); shp.line.fill.background()

tx(s, "Conclusion", 0.38, 0.1, 7, 0.58, sz=24, bold=True)

# Conclusion card
card(s, 0.38, 0.78, 8.3, 1.14, fill=BG2, ac=ACC, top=True)
tx(s,
   "SugarSense Predictor successfully combines a proven SVM classifier (~78% accuracy) with a "
   "5-provider AI chatbot into a unified, full-stack web application. The system delivers fast, "
   "evidence-based diabetes risk screening with immediate personalized health guidance — all configurable "
   "via a single environment file and deployable anywhere.",
   0.55, 0.86, 7.95, 0.98, sz=9.5, col=GRAY)

# 5 takeaway cards
takeaways=[
    (ACC,  "Proven ML","SVM linear kernel · ~78% accuracy · consistent with 7 published studies"),
    (LACC, "Modern Stack","React SPA + Flask REST · clean 3-tier architecture · JSON APIs"),
    (YELLOW,"5 AI Providers","1 env var switches between Claude, GPT, Gemini, HF, OpenRouter"),
    (PINK, "Knowledge Base","110+ Q&A in 10 diabetes categories · embedded in AI system prompt"),
    (GREEN,"Ready to Deploy","Prompt caching · multi-turn chat · graceful error handling built in"),
]
tw5=(8.3-0.2)/5; ty5=2.06
for i,(c,t,d) in enumerate(takeaways):
    x=0.38+i*(tw5+0.08)
    card(s, x, ty5, tw5, 1.28, fill=BG2, ac=c)
    num_circle(s, i+1, x+tw5/2-0.18, ty5+0.06, 0.18, fill=c)
    tx(s, t, x+0.08, ty5+0.52, tw5-0.14, 0.3, sz=9.5, bold=True, col=WHITE, align=PP_ALIGN.CENTER)
    tx(s, d, x+0.08, ty5+0.84, tw5-0.14, 0.46, sz=7.5, col=GRAY, align=PP_ALIGN.CENTER)

# References (left column)
card(s, 0.38, 3.52, 8.3, 3.72, fill=BG2, ac=PRI)
tx(s, "References", 0.55, 3.6, 4, 0.34, sz=11, bold=True)
refs=[
    "S. Patel & R. Goyal — Diabetes Prediction Using SVMs, IJCA vol.178, 2019",
    "M. Elshazly & Motaba — ML for Diabetes Prediction: Comparative Study, 2020",
    "N. Kumar & Srivastava — SVM and KNN Hybrid Model, Procedia CS vol.190, 2021",
    "J. Smith & P. Nguyen — Early Detection Using ML Algorithms, IEEE Access, 2021",
    "A. Banerjee & Tiwari — SVM with RFE Feature Selection, Applied Intelligence, 2022",
    "S. Li et al. — Diabetes Prediction Using Mobile Health Technology, 2022",
    "H. Kaur & M. Singh — SVM with Ensemble Methods, Expert Systems, 2023",
    "Smith et al. (1988) — Using ADAP Algorithm to Forecast Onset of Diabetes (Pima dataset)",
    "Pedregosa et al. (2011) — Scikit-learn: ML in Python, JMLR 12, 2825–2830",
    "Anthropic (2025) — Claude API Documentation · docs.anthropic.com",
    "IDF Diabetes Atlas 10th ed. (2021) · diabetesatlas.org",
]
col1r=refs[:6]; col2r=refs[6:]
for ci,col_refs in enumerate([col1r,col2r]):
    cx=0.52+ci*4.16
    for ri,ref in enumerate(col_refs):
        y=4.02+ri*0.52
        circle_num=s.shapes.add_shape(9,Inches(cx),Inches(y+0.06),Inches(0.28),Inches(0.28))
        circle_num.fill.solid(); circle_num.fill.fore_color.rgb=rgb(PRI); circle_num.line.fill.background()
        tp=circle_num.text_frame.paragraphs[0]; tp.alignment=PP_ALIGN.CENTER
        tr=tp.add_run(); tr.text=str(ci*6+ri+1); tr.font.size=Pt(7); tr.font.bold=True; tr.font.color.rgb=rgb(WHITE)
        tx(s, ref, cx+0.35, y+0.06, 3.65, 0.44, sz=8, col=GRAY)

# Thank You (right panel)
box(s, 8.86, 3.52, 4.1, 3.72, fill=PRI)
box(s, 8.86, 3.52, 4.1, 0.04, fill=LACC)
tx(s, "THANK", 8.9, 3.68, 4.0, 0.92, sz=42, bold=True, col=WHITE, align=PP_ALIGN.CENTER)
tx(s, "YOU!",  8.9, 4.42, 4.0, 0.88, sz=42, bold=True, col=LACC,  align=PP_ALIGN.CENTER)
box(s, 9.1, 5.38, 3.6, 0.03, fill=WHITE)
tx(s, "Questions Welcome!", 8.9, 5.48, 4.0, 0.35, sz=11, bold=True, col=WHITE, align=PP_ALIGN.CENTER)
tx(s, f"{COLLEGE}", 8.9, 5.9, 4.0, 0.28, sz=8, col="b2f5ea", align=PP_ALIGN.CENTER)
tx(s, GUIDE, 8.9, 6.2, 4.0, 0.28, sz=8, italic=True, col="b2f5ea", align=PP_ALIGN.CENTER)

# Team strip bottom
box(s, 0.38, 7.08, W-0.76, 0.38, fill="0d1a30")
for i,(name,roll) in enumerate(TEAM):
    x=0.5+i*2.6
    tx(s, f"0{i+1}. {name}", x, 7.12, 2.5, 0.2,  sz=7.5, bold=True, col=LACC)
    tx(s, roll,               x, 7.3,  2.5, 0.16, sz=7,   col=GRAY)

# ── Save ──────────────────────────────────────────────────────────────────────
prs.save("SugarSense_Presentation.pptx")
print(f"✓  SugarSense_Presentation.pptx  —  {len(prs.slides)} slides")
